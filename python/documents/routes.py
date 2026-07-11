"""
FastAPI routes for document generation: PowerPoint, Excel, and PDF.
"""

import os
from pathlib import Path

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel

from database import analytics_dao

router = APIRouter()


# ─── Models ───────────────────────────────────────────────────────────────────

class PresentationRequest(BaseModel):
    title: str
    slides: list[dict]  # [{"title": "...", "content": "...", "bullet_points": ["..."]}]
    template: str = "professional"

class SpreadsheetRequest(BaseModel):
    title: str
    sheets: list[dict]  # [{"name": "...", "headers": [...], "rows": [[...]]}]

class PDFRequest(BaseModel):
    title: str
    content: str  # Markdown content
    template: str = "professional"


# ─── PowerPoint Generation ────────────────────────────────────────────────────

@router.post("/powerpoint")
async def generate_presentation(request: PresentationRequest):
    """Generate a PowerPoint presentation from structured data."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = request.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(16, 185, 129)  # Emerald

        # Generate content slides
        for i, slide_data in enumerate(request.slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Slide title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", f"Slide {i+1}")
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(16, 185, 129)

            # Content
            content = slide_data.get("content", "")
            bullet_points = slide_data.get("bullet_points", [])

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True

            if content:
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
                p.space_after = Pt(12)

            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(16)
                p.level = 0
                p.space_after = Pt(6)

        # Save
        output_dir = Path(os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated"))
        output_dir.mkdir(exist_ok=True)
        output_path = output_dir / f"{request.title.replace(' ', '_')[:50]}.pptx"
        prs.save(str(output_path))

        await analytics_dao.log_activity(
            "documents", "generate_ppt", f"Generated presentation: {request.title}"
        )
        return {
            "status": "generated",
            "title": request.title,
            "slides": len(request.slides) + 1,
            "file_path": str(output_path),
            "file_size_bytes": output_path.stat().st_size,
        }
    except ImportError:
        return {"status": "unavailable", "message": "python-pptx not installed. Run: pip install python-pptx"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Excel Generation ─────────────────────────────────────────────────────────

@router.post("/excel")
async def generate_spreadsheet(request: SpreadsheetRequest):
    """Generate an Excel spreadsheet from structured data."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill

        wb = Workbook()

        for i, sheet_data in enumerate(request.sheets):
            if i == 0:
                ws = wb.active
            else:
                ws = wb.create_sheet()

            ws.title = sheet_data.get("name", f"Sheet{i+1}")
            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])

            # Header styling
            header_fill = PatternFill(start_color="10B981", end_color="10B981", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=12)

            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

            # Data rows
            for row_idx, row_data in enumerate(rows, 2):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.alignment = Alignment(horizontal="center")

            # Auto-fit column widths
            for col_idx in range(1, len(headers) + 1):
                max_length = max(
                    len(str(headers[col_idx - 1])),
                    max((len(str(row[col_idx - 1])) for row in rows if row), 0)
                )
                ws.column_dimensions[chr(64 + col_idx)].width = min(max_length + 2, 50)

        # Save
        output_dir = Path(os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated"))
        output_dir.mkdir(exist_ok=True)
        output_path = output_dir / f"{request.title.replace(' ', '_')[:50]}.xlsx"
        wb.save(str(output_path))

        await analytics_dao.log_activity(
            "documents", "generate_excel", f"Generated spreadsheet: {request.title}"
        )
        return {
            "status": "generated",
            "title": request.title,
            "sheets": len(request.sheets),
            "file_path": str(output_path),
            "file_size_bytes": output_path.stat().st_size,
        }
    except ImportError:
        return {"status": "unavailable", "message": "openpyxl not installed. Run: pip install openpyxl"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── PDF Generation ───────────────────────────────────────────────────────────

@router.post("/pdf")
async def generate_pdf(request: PDFRequest):
    """Generate a PDF document from markdown content."""
    try:
        # Try WeasyPrint first, fall back to reportlab
        try:
            import markdown
            from weasyprint import HTML

            html_content = markdown.markdown(request.content, extensions=["extra", "codehilite"])

            styled_html = f"""
            <html>
            <head>
                <style>
                    body {{ font-family: 'Georgia', serif; padding: 40px; color: #333; line-height: 1.6; }}
                    h1 {{ color: #10B981; border-bottom: 2px solid #10B981; padding-bottom: 8px; }}
                    h2 {{ color: #059669; margin-top: 24px; }}
                    code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }}
                    pre {{ background: #f4f4f4; padding: 12px; border-radius: 6px; }}
                    blockquote {{ border-left: 3px solid #10B981; padding-left: 16px; color: #666; }}
                </style>
            </head>
            <body>
                <h1>{request.title}</h1>
                {html_content}
            </body>
            </html>
            """

            output_dir = Path(os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated"))
            output_dir.mkdir(exist_ok=True)
            output_path = output_dir / f"{request.title.replace(' ', '_')[:50]}.pdf"

            HTML(string=styled_html).write_pdf(str(output_path))

        except ImportError:
            # Fallback to reportlab
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

            output_dir = Path(os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated"))
            output_dir.mkdir(exist_ok=True)
            output_path = output_dir / f"{request.title.replace(' ', '_')[:50]}.pdf"

            doc = SimpleDocTemplate(str(output_path), pagesize=letter)
            styles = getSampleStyleSheet()

            story = [
                Paragraph(request.title, styles["Title"]),
                Spacer(1, 12),
                Paragraph(request.content[:5000].replace("\n", "<br/>"), styles["Normal"]),
            ]
            doc.build(story)

        await analytics_dao.log_activity(
            "documents", "generate_pdf", f"Generated PDF: {request.title}"
        )
        return {
            "status": "generated",
            "title": request.title,
            "file_path": str(output_path),
            "file_size_bytes": output_path.stat().st_size,
        }
    except ImportError:
        return {"status": "unavailable",
                "message": "PDF libraries not installed. Run: pip install weasyprint markdown (or reportlab as fallback)"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ─── Export Jobs to CSV (quick utility) ───────────────────────────────────────

@router.get("/export/jobs")
async def export_jobs_csv():
    """Export job listings to CSV."""
    try:
        from database import jobs_dao

        jobs = await jobs_dao.get_active_jobs()
        if not jobs:
            return {"status": "no_data", "message": "No jobs to export"}

        import csv
        output_dir = Path(os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated"))
        output_dir.mkdir(exist_ok=True)
        output_path = output_dir / "jobs_export.csv"

        with open(output_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(["Title", "Company", "Location", "Salary Min", "Salary Max", "Source", "URL"])
            for job in jobs:
                writer.writerow([
                    job.get("title", ""), job.get("company", ""),
                    job.get("location", ""), job.get("salary_min", 0),
                    job.get("salary_max", 0), job.get("source_board", ""),
                    job.get("source_url", ""),
                ])

        return {"status": "exported", "file_path": str(output_path), "count": len(jobs)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
