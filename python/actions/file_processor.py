"""
file_processor.py — BARQ Universal File Processor

Supported types:
  image   → describe, ocr, resize, convert, compress, crop, info
  pdf     → summarize, extract_text, info, to_word
  docx    → summarize, extract_text, reformat, word_count
  txt/md  → summarize, reformat, word_count
  csv     → analyze, info, stats, filter, sort, convert
  xlsx    → analyze, info, stats, filter, sort, convert
  json    → validate, format, analyze, to_csv
  code    → explain, review, fix, optimize, document, run, info
  audio   → transcribe, info, convert, trim
  video   → info, extract_audio, trim, extract_frame, compress, transcribe, convert
  zip/tar → list, extract
  pptx    → summarize, extract_text, analyze

Requires optional dependencies:
  Pillow, pdfplumber (or PyPDF2), python-docx, pandas, openpyxl,
  pydub, python-pptx, ffmpeg (system)
"""

import json
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Callable

# ─── Helpers ───────────────────────────────────────────────────────────

def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "svg", "ico"}
    video_exts = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}
    audio_exts = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "wma", "opus"}
    code_exts  = {"py", "js", "ts", "jsx", "tsx", "html", "css", "java", "c",
                  "cpp", "cs", "go", "rs", "rb", "php", "swift", "kt", "sh",
                  "bash", "ps1", "lua", "r", "m", "sql", "yaml", "toml"}
    archive_exts = {"zip", "rar", "tar", "gz", "7z", "bz2", "xz"}

    if ext in image_exts:   return "image"
    if ext in video_exts:   return "video"
    if ext in audio_exts:   return "audio"
    if ext in code_exts:    return "code"
    if ext in archive_exts: return "archive"
    if ext == "pdf":        return "pdf"
    if ext in ("docx", "doc"): return "docx"
    if ext in ("txt", "md", "rst", "log"): return "text"
    if ext in ("csv", "tsv"): return "csv"
    if ext in ("xlsx", "xls", "ods"): return "excel"
    if ext == "json":       return "json"
    if ext == "xml":        return "xml"
    if ext in ("pptx", "ppt"): return "pptx"
    return "unknown"


def _file_size_str(path: Path) -> str:
    size = path.stat().st_size
    if size < 1024:        return f"{size} B"
    if size < 1024**2:     return f"{size/1024:.1f} KB"
    if size < 1024**3:     return f"{size/1024**2:.1f} MB"
    return f"{size/1024**3:.1f} GB"


def _output_path(src: Path, suffix: str, new_ext: str = None) -> Path:
    ext = new_ext or src.suffix
    name = f"{src.stem}_{suffix}{ext}"
    return src.parent / name


def _gemini_text(prompt: str) -> str:
    """Run a text-only Gemini prompt and return the response."""
    import asyncio
    from agent.vision import analyze_image_with_gemini
    # Use analyze_image_with_gemini with a tiny placeholder image to get text-only response
    # Alternatively, use the genai client directly
    try:
        from google import genai
        from config import get_settings
        settings = get_settings()
        client = genai.Client(api_key=settings.gemini_api_key)
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt,
        )
        return response.text.strip()
    except Exception as e:
        return f"[Gemini error: {e}]"


def _run_async(coro):
    """Run a coroutine synchronously from a thread pool context."""
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(coro)
        loop.close()
        return result
    except RuntimeError:
        return asyncio.run(coro)


# ─── Image Processing ─────────────────────────────────────────────────

def _process_image(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "describe"
    try:
        from PIL import Image
    except ImportError:
        return {"status": "error", "detail": "Pillow not installed. Install: pip install Pillow"}

    # ── AI analysis (describe, ocr, analyze, read, extract_text) ─────
    if action in ("describe", "ocr", "analyze", "read", "extract_text"):
        try:
            img = Image.open(path)
            prompt_map = {
                "describe": "Describe this image in detail — objects, setting, colors, composition.",
                "ocr": "Extract all text visible in this image. Return only the text, preserving structure.",
                "analyze": "Analyze this image: objects, colors, composition, any text visible, context.",
                "read": "Read all text visible in this image.",
                "extract_text": "Extract all text from this image.",
            }
            prompt = prompt_map.get(action, "Describe this image.")
            if params.get("instruction"):
                prompt = params["instruction"]

            from agent.vision import analyze_image_with_gemini

            # Save image to bytes
            import io
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            image_bytes = buf.getvalue()

            result = _run_async(analyze_image_with_gemini(
                image_bytes, "image/png", prompt=prompt
            ))
            result = result.strip()

            if len(result) > 500 and params.get("save", True):
                out = _output_path(path, "result", ".txt")
                out.write_text(result, encoding="utf-8")
                return {
                    "status": "success",
                    "detail": f"{result[:300]}...\n\nFull result saved to: {out.name}",
                    "saved_to": str(out),
                }
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"AI image analysis failed: {e}"}

    # ── resize ─────────────────────────────────────────────────────
    if action == "resize":
        width = int(params.get("width", 0))
        height = int(params.get("height", 0))
        scale = float(params.get("scale", 0))
        try:
            img = Image.open(path)
            w, h = img.size
            if scale:
                new_size = (int(w * scale), int(h * scale))
            elif width and height:
                new_size = (width, height)
            elif width:
                new_size = (width, int(h * width / w))
            elif height:
                new_size = (int(w * height / h), height)
            else:
                return {"status": "error", "detail": "Specify width, height, or scale"}
            out = _output_path(path, f"resized_{new_size[0]}x{new_size[1]}")
            img.resize(new_size, Image.LANCZOS).save(out)
            return {"status": "success", "detail": f"Resized {w}x{h} -> {new_size[0]}x{new_size[1]}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Resize failed: {e}"}

    # ── convert ────────────────────────────────────────────────────
    if action == "convert":
        fmt = params.get("format", "png").lower().strip(".")
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG", "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}
        pil_fmt = fmt_map.get(fmt, fmt.upper())
        try:
            img = Image.open(path).convert("RGB") if fmt in ("jpg", "jpeg") else Image.open(path)
            out = _output_path(path, "converted", f".{fmt}")
            img.save(out, pil_fmt)
            return {"status": "success", "detail": f"Converted to {fmt.upper()}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Convert failed: {e}"}

    # ── compress ───────────────────────────────────────────────────
    if action == "compress":
        quality = int(params.get("quality", 70))
        try:
            img = Image.open(path).convert("RGB")
            out = _output_path(path, f"compressed_q{quality}", ".jpg")
            img.save(out, "JPEG", quality=quality, optimize=True)
            return {
                "status": "success",
                "detail": f"Compressed: {_file_size_str(path)} -> {_file_size_str(out)}. Saved: {out.name}",
                "saved_to": str(out),
            }
        except Exception as e:
            return {"status": "error", "detail": f"Compress failed: {e}"}

    # ── info ───────────────────────────────────────────────────────
    if action == "info":
        try:
            img = Image.open(path)
            return {
                "status": "success",
                "detail": f"Image: {img.format}, {img.size[0]}x{img.size[1]}px, mode: {img.mode}, size: {_file_size_str(path)}",
                "format": img.format,
                "width": img.size[0],
                "height": img.size[1],
                "mode": img.mode,
            }
        except Exception as e:
            return {"status": "error", "detail": f"Info failed: {e}"}

    return {"status": "error", "detail": f"Unknown image action: '{action}'. Try: describe, ocr, resize, convert, compress, info"}


# ─── PDF Processing ───────────────────────────────────────────────────

def _process_pdf(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "summarize"

    def _extract_pdf_text(max_chars=50000) -> str:
        text = ""
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text += (page.extract_text() or "") + "\n"
        except ImportError:
            try:
                import PyPDF2
                with open(path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            except ImportError:
                pass
        return text[:max_chars]

    if action in ("summarize", "extract_text", "analyze", "reformat"):
        text = _extract_pdf_text()
        if not text.strip():
            return {"status": "error", "detail": "Could not extract text from PDF (may be scanned/image-based). Try OCR."}

        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(text, encoding="utf-8")
            return {"status": "success", "detail": f"Text extracted ({len(text)} chars). Saved: {out.name}", "saved_to": str(out)}

        prompt_map = {
            "summarize": f"Summarize this PDF document concisely:\n\n{text}",
            "analyze": f"Analyze this document thoroughly:\n\n{text}",
            "reformat": f"Reformat this text cleanly:\n\n{text}",
        }
        try:
            result = _gemini_text(prompt_map.get(action, f"Analyze:\n\n{text}"))
            if len(result) > 600 and params.get("save", True):
                out = _output_path(path, action, ".txt")
                out.write_text(result, encoding="utf-8")
                return {"status": "success", "detail": f"{result[:400]}...\n\nFull result saved: {out.name}", "saved_to": str(out)}
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"AI analysis failed: {e}"}

    if action == "info":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
            return {"status": "success", "detail": f"PDF: {pages} pages, size: {_file_size_str(path)}", "pages": pages}
        except Exception:
            return {"status": "success", "detail": f"PDF size: {_file_size_str(path)}"}

    if action == "to_word":
        text = _extract_pdf_text()
        if not text:
            return {"status": "error", "detail": "Could not extract text to convert."}
        try:
            from docx import Document
            doc = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "converted", ".docx")
            doc.save(out)
            return {"status": "success", "detail": f"Converted to Word. Saved: {out.name}", "saved_to": str(out)}
        except ImportError:
            return {"status": "error", "detail": "python-docx not installed. Install: pip install python-docx"}

    return {"status": "error", "detail": f"Unknown PDF action: '{action}'. Try: summarize, extract_text, info, to_word"}


# ─── Document Processing (docx, txt, md) ──────────────────────────────

def _process_text_doc(path: Path, file_type: str, action: str, params: dict) -> dict[str, Any]:
    action = action or "summarize"

    def _read_content() -> str:
        if file_type == "docx":
            try:
                from docx import Document
                doc = Document(path)
                return "\n".join(p.text for p in doc.paragraphs)
            except ImportError:
                return ""
            except Exception as e:
                return f"[Read error: {e}]"
        return path.read_text(encoding="utf-8", errors="ignore")

    content = _read_content()
    if not content.strip():
        return {"status": "error", "detail": "File appears to be empty or could not be read."}

    if action == "word_count":
        words = len(content.split())
        chars = len(content)
        lines = content.count("\n")
        return {"status": "success", "detail": f"Words: {words}, Characters: {chars}, Lines: {lines}", "words": words, "chars": chars, "lines": lines}

    if action == "extract_text":
        if file_type != "txt":
            out = _output_path(path, "extracted", ".txt")
            out.write_text(content, encoding="utf-8")
            return {"status": "success", "detail": f"Text extracted. Saved: {out.name}", "saved_to": str(out)}
        return {"status": "success", "detail": content[:2000]}

    prompt_map = {
        "summarize": f"Summarize this document concisely:\n\n{content[:40000]}",
        "analyze": f"Analyze this document:\n\n{content[:40000]}",
        "reformat": f"Reformat this text with clean structure:\n\n{content[:40000]}",
        "to_bullet": f"Convert to bullet-point summary:\n\n{content[:40000]}",
    }
    instruction = params.get("instruction", "")
    if action not in prompt_map:
        prompt = f"{instruction or action}\n\n{content[:40000]}"
    else:
        prompt = prompt_map[action]

    try:
        result = _gemini_text(prompt)
        if len(result) > 600 and params.get("save", True):
            out = _output_path(path, action, ".txt")
            out.write_text(result, encoding="utf-8")
            return {"status": "success", "detail": f"{result[:400]}...\n\nFull result saved: {out.name}", "saved_to": str(out)}
        return {"status": "success", "detail": result}
    except Exception as e:
        return {"status": "error", "detail": f"AI processing failed: {e}"}


# ─── Data Processing (CSV, Excel) ─────────────────────────────────────

def _process_data(path: Path, file_type: str, action: str, params: dict) -> dict[str, Any]:
    try:
        import pandas as pd
    except ImportError:
        return {"status": "error", "detail": "pandas not installed. Install: pip install pandas openpyxl"}

    action = action or "analyze"
    try:
        if file_type == "csv":
            df = pd.read_csv(path, encoding="utf-8", errors="replace")
        else:
            df = pd.read_excel(path)
    except Exception as e:
        return {"status": "error", "detail": f"Could not read file: {e}"}

    if action == "info":
        return {
            "status": "success",
            "detail": f"Rows: {len(df)}, Columns: {len(df.columns)}\nColumns: {', '.join(df.columns.tolist())}\nSize: {_file_size_str(path)}",
            "rows": len(df), "columns": len(df.columns), "column_names": df.columns.tolist(),
        }

    if action == "stats":
        try:
            desc = df.describe(include="all").to_string()
            return {"status": "success", "detail": f"Statistics:\n{desc[:2000]}"}
        except Exception as e:
            return {"status": "error", "detail": f"Stats failed: {e}"}

    if action == "analyze":
        preview = df.head(50).to_string()
        prompt = f"Analyze this dataset. Columns: {list(df.columns)}\nRows: {len(df)}\nPreview:\n{preview}\n\nGive insights, patterns, and notable findings."
        try:
            result = _gemini_text(prompt)
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"AI analysis failed: {e}"}

    if action in ("convert", "to_csv", "to_excel", "to_json"):
        fmt = {"to_csv": "csv", "to_excel": "xlsx", "to_json": "json", "convert": params.get("format", "csv")}.get(action, "csv")
        try:
            if fmt == "csv":
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False, encoding="utf-8")
            elif fmt == "xlsx":
                out = _output_path(path, "converted", ".xlsx")
                df.to_excel(out, index=False)
            elif fmt == "json":
                out = _output_path(path, "converted", ".json")
                df.to_json(out, orient="records", force_ascii=False, indent=2)
            return {"status": "success", "detail": f"Converted to {fmt.upper()}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Convert failed: {e}"}

    if action == "filter":
        col = params.get("column", "")
        value = params.get("value", "")
        condition = params.get("condition", "equals")
        if not col or col not in df.columns:
            return {"status": "error", "detail": f"Column '{col}' not found. Available: {', '.join(df.columns)}"}
        try:
            if condition == "equals":     filtered = df[df[col] == value]
            elif condition == "contains": filtered = df[df[col].astype(str).str.contains(str(value), case=False)]
            elif condition == "gt":       filtered = df[df[col] > float(value)]
            elif condition == "lt":       filtered = df[df[col] < float(value)]
            else:                         filtered = df[df[col] == value]
            out = _output_path(path, "filtered", ".csv")
            filtered.to_csv(out, index=False)
            return {"status": "success", "detail": f"Filtered: {len(filtered)} rows match. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Filter failed: {e}"}

    if action == "sort":
        col = params.get("column", df.columns[0])
        asc = params.get("ascending", True)
        try:
            sorted_df = df.sort_values(col, ascending=asc)
            out = _output_path(path, "sorted", path.suffix)
            sorted_df.to_csv(out, index=False)
            return {"status": "success", "detail": f"Sorted by '{col}'. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Sort failed: {e}"}

    return {"status": "error", "detail": f"Unknown action: '{action}'. Try: analyze, info, stats, filter, sort, convert"}


# ─── JSON Processing ──────────────────────────────────────────────────

def _process_json(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "analyze"
    try:
        content = path.read_text(encoding="utf-8")
        data = json.loads(content)
    except Exception as e:
        return {"status": "error", "detail": f"Invalid JSON: {e}"}

    if action == "validate":
        return {"status": "success", "detail": f"Valid JSON. Type: {type(data).__name__}, size: {_file_size_str(path)}", "json_type": type(data).__name__}

    if action == "format":
        out = _output_path(path, "formatted", ".json")
        out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return {"status": "success", "detail": f"Formatted JSON saved: {out.name}", "saved_to": str(out)}

    if action in ("analyze", "summarize", "extract"):
        preview = json.dumps(data, indent=2, ensure_ascii=False)[:8000]
        instruction = params.get("instruction", "")
        prompt = f"{instruction}\n\nJSON data:\n{preview}" if instruction else f"Task: {action} this JSON data:\n{preview}"
        try:
            result = _gemini_text(prompt)
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"AI processing failed: {e}"}

    if action == "to_csv":
        try:
            import pandas as pd
            if isinstance(data, list):
                df = pd.DataFrame(data)
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False)
                return {"status": "success", "detail": f"Converted to CSV. Saved: {out.name}", "saved_to": str(out)}
            return {"status": "error", "detail": "JSON must be an array of objects to convert to CSV."}
        except ImportError:
            return {"status": "error", "detail": "pandas not installed."}

    return {"status": "error", "detail": f"Unknown JSON action: '{action}'. Try: validate, format, analyze, to_csv"}


# ─── Code Processing ──────────────────────────────────────────────────

def _process_code(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "explain"
    content = path.read_text(encoding="utf-8", errors="ignore")
    ext = path.suffix.lstrip(".")

    if action == "run":
        if ext == "py":
            try:
                result = subprocess.run(
                    ["python", str(path)], capture_output=True, text=True, timeout=30
                )
                out = result.stdout or result.stderr
                return {"status": "success", "detail": f"Output:\n{out[:2000]}" if out else "No output.", "return_code": result.returncode}
            except subprocess.TimeoutExpired:
                return {"status": "error", "detail": "Execution timed out (30s)."}
            except Exception as e:
                return {"status": "error", "detail": f"Run failed: {e}"}
        return {"status": "error", "detail": f"Direct execution not supported for .{ext} files."}

    if action == "info":
        lines = content.count("\n")
        words = len(content.split())
        return {"status": "success", "detail": f"Code: {lines} lines, {words} words, {_file_size_str(path)}", "lines": lines, "words": words}

    prompt_map = {
        "explain":   f"Explain this {ext} code clearly:\n\n```{ext}\n{content[:30000]}\n```",
        "review":    f"Review this {ext} code for bugs and improvements:\n\n```{ext}\n{content[:30000]}\n```",
        "fix":       f"Fix any bugs in this {ext} code, return corrected version:\n\n```{ext}\n{content[:30000]}\n```",
        "optimize":  f"Optimize this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
        "document":  f"Add documentation to this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
        "summarize": f"Summarize what this {ext} code does:\n\n```{ext}\n{content[:30000]}\n```",
    }
    instruction = params.get("instruction", "")
    prompt = prompt_map.get(action, f"{instruction or action}\n\n```{ext}\n{content[:30000]}\n```")

    try:
        result = _gemini_text(prompt)
        if action in ("fix", "optimize", "document") and params.get("save", True):
            out = _output_path(path, action)
            code_match = re.search(r"```(?:\w+)?\n(.*?)```", result, re.DOTALL)
            code_to_save = code_match.group(1) if code_match else result
            out.write_text(code_to_save, encoding="utf-8")
            return {"status": "success", "detail": f"{result[:400]}...\n\nSaved: {out.name}", "saved_to": str(out)}
        return {"status": "success", "detail": result}
    except Exception as e:
        return {"status": "error", "detail": f"AI processing failed: {e}"}


# ─── Audio Processing ─────────────────────────────────────────────────

def _process_audio(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "info"

    if action == "info":
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            duration = len(audio) / 1000
            mins, secs = divmod(int(duration), 60)
            return {
                "status": "success",
                "detail": f"Audio: {mins}m {secs}s, {audio.channels}ch, {audio.frame_rate}Hz, {_file_size_str(path)}",
                "duration_s": int(duration),
                "channels": audio.channels,
                "sample_rate": audio.frame_rate,
            }
        except ImportError:
            return {"status": "success", "detail": f"Audio file: {_file_size_str(path)} (install pydub for details)"}
        except Exception as e:
            return {"status": "error", "detail": f"Info failed: {e}"}

    if action == "transcribe":
        try:
            import base64
            mime_map = {"mp3": "audio/mp3", "wav": "audio/wav", "ogg": "audio/ogg", "m4a": "audio/mp4", "aac": "audio/aac", "flac": "audio/flac"}
            mime = mime_map.get(path.suffix.lstrip(".").lower(), "audio/mpeg")
            file_bytes = path.read_bytes()

            prompt = "Transcribe all speech in this audio file accurately."
            from google import genai
            from config import get_settings
            settings = get_settings()
            client = genai.Client(api_key=settings.gemini_api_key)
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[prompt, {"mime_type": mime, "data": file_bytes}],
            )
            result = response.text.strip()
            if params.get("save", True):
                out = _output_path(path, "transcript", ".txt")
                out.write_text(result, encoding="utf-8")
                return {"status": "success", "detail": f"Transcription saved: {out.name}\n\nPreview: {result[:300]}", "saved_to": str(out)}
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"Transcription failed: {e}"}

    if action == "convert":
        fmt = params.get("format", "mp3").lstrip(".")
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            out = _output_path(path, "converted", f".{fmt}")
            audio.export(out, format=fmt)
            return {"status": "success", "detail": f"Converted to {fmt.upper()}. Saved: {out.name}", "saved_to": str(out)}
        except ImportError:
            return {"status": "error", "detail": "pydub not installed. Install: pip install pydub"}
        except Exception as e:
            return {"status": "error", "detail": f"Convert failed: {e}"}

    if action == "trim":
        start = float(params.get("start", 0))
        end = float(params.get("end", 0))
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            end_ms = int(end * 1000) if end else len(audio)
            trimmed = audio[int(start * 1000):end_ms]
            out = _output_path(path, f"trim_{int(start)}s_{int(end)}s")
            trimmed.export(out, format=path.suffix.lstrip("."))
            return {"status": "success", "detail": f"Trimmed {int(start)}s-{int(end)}s. Saved: {out.name}", "saved_to": str(out)}
        except ImportError:
            return {"status": "error", "detail": "pydub not installed."}
        except Exception as e:
            return {"status": "error", "detail": f"Trim failed: {e}"}

    return {"status": "error", "detail": f"Unknown audio action: '{action}'. Try: transcribe, info, convert, trim"}


# ─── Video Processing ─────────────────────────────────────────────────

def _process_video(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "info"

    def _ffmpeg_available() -> bool:
        try:
            subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=3)
            return True
        except Exception:
            return False

    if action == "info":
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", "-show_streams", str(path)],
                capture_output=True, text=True, timeout=10
            )
            data = json.loads(result.stdout)
            fmt = data.get("format", {})
            duration = float(fmt.get("duration", 0))
            mins, secs = divmod(int(duration), 60)
            streams = data.get("streams", [])
            video_s = next((s for s in streams if s["codec_type"] == "video"), {})
            w, h = video_s.get("width", "?"), video_s.get("height", "?")
            fps = video_s.get("r_frame_rate", "?")
            return {"status": "success", "detail": f"Video: {mins}m {secs}s, {w}x{h}, {fps} fps, {_file_size_str(path)}"}
        except Exception:
            return {"status": "success", "detail": f"Video file: {_file_size_str(path)}"}

    if not _ffmpeg_available():
        return {"status": "error", "detail": "ffmpeg not found. Install ffmpeg for video processing."}

    if action == "extract_audio":
        out = _output_path(path, "audio", ".mp3")
        try:
            subprocess.run(["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a", str(out), "-y"], capture_output=True, timeout=300)
            return {"status": "success", "detail": f"Audio extracted. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Extract audio failed: {e}"}

    if action == "trim":
        start = params.get("start", "00:00:00")
        end = params.get("end", "")
        out = _output_path(path, "trim", path.suffix)
        try:
            cmd = ["ffmpeg", "-i", str(path), "-ss", str(start)]
            if end:
                cmd += ["-to", str(end)]
            cmd += ["-c", "copy", str(out), "-y"]
            subprocess.run(cmd, capture_output=True, timeout=600)
            return {"status": "success", "detail": f"Trimmed video saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Trim failed: {e}"}

    if action == "extract_frame":
        timestamp = params.get("timestamp", "00:00:01")
        out = _output_path(path, f"frame_{timestamp.replace(':', '')}", ".jpg")
        try:
            subprocess.run(["ffmpeg", "-i", str(path), "-ss", timestamp, "-vframes", "1", str(out), "-y"], capture_output=True, timeout=30)
            return {"status": "success", "detail": f"Frame at {timestamp}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Extract frame failed: {e}"}

    if action == "compress":
        crf = int(params.get("quality", 28))
        out = _output_path(path, f"compressed_crf{crf}", ".mp4")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-c:v", "libx264", "-crf", str(crf), "-preset", "medium", "-c:a", "copy", str(out), "-y"],
                capture_output=True, timeout=1800
            )
            return {"status": "success", "detail": f"Compressed: {_file_size_str(path)} -> {_file_size_str(out)}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Compress failed: {e}"}

    if action == "transcribe":
        tmp_audio = Path(tempfile.mktemp(suffix=".mp3"))
        try:
            subprocess.run(["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a", str(tmp_audio), "-y"], capture_output=True, timeout=300)
            return _process_audio(tmp_audio, "transcribe", params)
        except Exception as e:
            return {"status": "error", "detail": f"Video transcription failed: {e}"}
        finally:
            if tmp_audio.exists():
                tmp_audio.unlink()

    if action == "convert":
        fmt = params.get("format", "mp4").lstrip(".")
        out = _output_path(path, "converted", f".{fmt}")
        try:
            subprocess.run(["ffmpeg", "-i", str(path), str(out), "-y"], capture_output=True, timeout=1800)
            return {"status": "success", "detail": f"Converted to {fmt.upper()}. Saved: {out.name}", "saved_to": str(out)}
        except Exception as e:
            return {"status": "error", "detail": f"Convert failed: {e}"}

    return {"status": "error", "detail": f"Unknown video action: '{action}'. Try: info, extract_audio, trim, extract_frame, compress, transcribe, convert"}


# ─── Archive Processing ───────────────────────────────────────────────

def _process_archive(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "list"

    if action == "list":
        try:
            import zipfile, tarfile
            ext = path.suffix.lower()
            if ext == ".zip":
                with zipfile.ZipFile(path) as z:
                    names = z.namelist()
            elif ext in (".tar", ".gz", ".bz2", ".xz"):
                with tarfile.open(path) as t:
                    names = t.getnames()
            else:
                return {"status": "error", "detail": f"Unsupported archive: {ext}"}
            preview = "\n".join(names[:30])
            suffix = f"\n... and {len(names)-30} more" if len(names) > 30 else ""
            return {"status": "success", "detail": f"Archive: {len(names)} files\n{preview}{suffix}", "file_count": len(names)}
        except Exception as e:
            return {"status": "error", "detail": f"List failed: {e}"}

    if action == "extract":
        dest = Path(params.get("destination", str(path.parent / path.stem)))
        dest.mkdir(parents=True, exist_ok=True)
        try:
            shutil.unpack_archive(path, dest)
            return {"status": "success", "detail": f"Extracted to: {dest}", "saved_to": str(dest)}
        except Exception as e:
            return {"status": "error", "detail": f"Extract failed: {e}"}

    return {"status": "error", "detail": f"Unknown archive action: '{action}'. Try: list, extract"}


# ─── PPTX Processing ──────────────────────────────────────────────────

def _process_pptx(path: Path, action: str, params: dict) -> dict[str, Any]:
    action = action or "summarize"

    def _read_pptx_text() -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {i} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                text.append(slide_text)
            return "\n".join(text)
        except ImportError:
            return ""
        except Exception as e:
            return f"[Read error: {e}]"

    if action in ("summarize", "extract_text", "analyze"):
        text = _read_pptx_text()
        if not text:
            return {"status": "error", "detail": "python-pptx not installed or file empty. Install: pip install python-pptx"}

        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(text, encoding="utf-8")
            return {"status": "success", "detail": f"Text extracted. Saved: {out.name}", "saved_to": str(out)}

        try:
            prompt = f"{'Summarize' if action == 'summarize' else 'Analyze'} this presentation:\n{text[:30000]}"
            result = _gemini_text(prompt)
            return {"status": "success", "detail": result}
        except Exception as e:
            return {"status": "error", "detail": f"AI processing failed: {e}"}

    return {"status": "error", "detail": f"Unknown PPTX action: '{action}'. Try: summarize, extract_text, analyze"}


# ─── Unknown File Type ────────────────────────────────────────────────

def _process_unknown(path: Path, action: str, params: dict) -> dict[str, Any]:
    try:
        content = path.read_text(encoding="utf-8", errors="ignore")[:10000]
        instruction = params.get("instruction", action or "Describe what this file contains and what can be done with it.")
        prompt = f"File: {path.name}\nContent preview:\n{content}\n\nTask: {instruction}"
        result = _gemini_text(prompt)
        return {"status": "success", "detail": result}
    except Exception as e:
        return {"status": "error", "detail": f"Unknown file type ({path.suffix}). Could not process: {e}"}


# ─── Public API ───────────────────────────────────────────────────────

def process_file(
    file_path: str = "",
    action: str = "",
    instruction: str = "",
    save: bool = True,
    **kwargs,
) -> dict[str, Any]:
    """Process a file with AI-powered analysis and transformations.

    Detects file type automatically and dispatches to the appropriate handler.

    Args:
        file_path: Absolute or relative path to the file (required).
        action: What to do with the file (depends on type).
                Image: describe, ocr, resize, convert, compress, info
                PDF: summarize, extract_text, info, to_word
                Docx/Text: summarize, extract_text, word_count, reformat
                CSV/Excel: analyze, info, stats, filter, sort, convert
                JSON: validate, format, analyze, to_csv
                Code: explain, review, fix, optimize, document, run, info
                Audio: transcribe, info, convert, trim
                Video: info, extract_audio, trim, extract_frame, compress, transcribe, convert
                Archive: list, extract
                PPTX: summarize, extract_text, analyze
        instruction: Custom instruction override for AI analysis.
        save: Whether to save long results to disk (default: True).
        **kwargs: Additional params (width, height, quality, format, etc.).

    Returns:
        Dict with status and detail.
    """
    if not file_path:
        return {"status": "error", "detail": "No file path provided."}

    path = Path(file_path)
    if not path.exists():
        return {"status": "error", "detail": f"File not found: {file_path}"}
    if not path.is_file():
        return {"status": "error", "detail": f"Path is not a file: {file_path}"}

    file_type = _detect_type(path)
    params = {**kwargs, "save": save}
    if instruction:
        params["instruction"] = instruction

    print(f"[FileProcessor] {file_type.upper()} | {path.name} | action={action or 'auto'}")

    dispatch: dict[str, Callable] = {
        "image":   _process_image,
        "pdf":     _process_pdf,
        "docx":    lambda p, a, pm: _process_text_doc(p, "docx", a, pm),
        "text":    lambda p, a, pm: _process_text_doc(p, "text", a, pm),
        "csv":     lambda p, a, pm: _process_data(p, "csv", a, pm),
        "excel":   lambda p, a, pm: _process_data(p, "excel", a, pm),
        "json":    _process_json,
        "xml":     lambda p, a, pm: _process_json(p, a, pm),
        "code":    _process_code,
        "audio":   _process_audio,
        "video":   _process_video,
        "archive": _process_archive,
        "pptx":    _process_pptx,
        "unknown": _process_unknown,
    }

    handler = dispatch.get(file_type)
    if not handler:
        return {"status": "error", "detail": f"Unsupported file type: {file_type}"}

    try:
        result = handler(path, action, params)
        return result
    except Exception as e:
        import traceback
        traceback.print_exc()
        return {"status": "error", "detail": f"Processing failed: {e}"}
